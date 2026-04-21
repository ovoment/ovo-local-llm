"""Document parsing dispatcher — routes files through kordoc CLI subprocess.

Usage:
    doc = await parse_file(Path("/path/to/report.pdf"))
    print(doc.full_text)   # Markdown
    print(doc.pages)       # page count
"""
from __future__ import annotations

import asyncio
import json
import logging
import mimetypes
import subprocess
from pathlib import Path

from ovo_sidecar.parsing.kordoc_installer import is_ready, node_bin, kordoc_cli
from ovo_sidecar.parsing.models import (
    ParsedDocument,
    ParsedSection,
    ParsedTable,
    file_hash,
    _new_id,
    _now_kst,
)

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".hwp", ".hwpx",
    ".txt", ".md", ".csv", ".json",
    ".xml",
}

MIME_MAP: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".hwp": "application/x-hwp",
    ".hwpx": "application/x-hwpx",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".csv": "text/csv",
    ".json": "application/json",
    ".xml": "application/xml",
}


def _estimate_tokens(text: str) -> int:
    """Rough token estimate: ~1.3 tokens per Korean char, ~0.75 per English word."""
    if not text:
        return 0
    ascii_chars = sum(1 for c in text if ord(c) < 128)
    non_ascii = len(text) - ascii_chars
    return int(non_ascii * 1.3 + ascii_chars * 0.25)


def _detect_mime(path: Path) -> str:
    ext = path.suffix.lower()
    return MIME_MAP.get(ext, mimetypes.guess_type(str(path))[0] or "application/octet-stream")


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


async def parse_file(path: Path) -> ParsedDocument:
    """Parse a document file using kordoc CLI. Raises RuntimeError if kordoc not installed."""
    if not is_ready():
        raise RuntimeError("kordoc not installed — call /ovo/parse/install first")

    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()

    if ext in (".txt", ".md", ".csv", ".json"):
        return await _parse_plaintext(path)

    if ext == ".pptx":
        return await _parse_pptx(path)

    doc = await _parse_with_kordoc(path)

    if ext == ".pdf" and len(doc.full_text.strip()) < 50 and doc.pages > 0:
        logger.info("PDF appears to be scanned/image-based, attempting OCR: %s", path.name)
        try:
            ocr_doc = await _ocr_pdf(path)
            if len(ocr_doc.full_text.strip()) > len(doc.full_text.strip()):
                return ocr_doc
        except Exception as e:
            logger.warning("OCR failed for %s: %s", path.name, e)

    return doc


async def _parse_pptx(path: Path) -> ParsedDocument:
    """Parse PPTX using python-pptx (kordoc doesn't support PPTX)."""
    def _extract() -> tuple[str, list[ParsedSection]]:
        from pptx import Presentation  # type: ignore[import-untyped]
        prs = Presentation(str(path))
        parts: list[str] = []
        sections: list[ParsedSection] = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            if not title:
                                title = text
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(row_texts))
            slide_md = f"# Slide {i}" + (f": {title}" if title else "") + "\n\n" + "\n".join(texts)
            parts.append(slide_md)
            if title:
                sections.append(ParsedSection(title=title, heading_level=1, page=i))
        return "\n\n---\n\n".join(parts), sections

    full_text, sections = await asyncio.to_thread(_extract)
    return ParsedDocument(
        doc_id=_new_id(),
        filename=path.name,
        mime=_detect_mime(path),
        source_path=str(path),
        file_hash=file_hash(path),
        pages=len(sections),
        full_text=full_text,
        sections=sections,
        tokens_estimate=_estimate_tokens(full_text),
        parsed_at=_now_kst(),
    )


async def _parse_plaintext(path: Path) -> ParsedDocument:
    """Direct read for plaintext formats — no kordoc needed."""
    def _read() -> str:
        return path.read_text(encoding="utf-8", errors="replace")

    text = await asyncio.to_thread(_read)

    return ParsedDocument(
        doc_id=_new_id(),
        filename=path.name,
        mime=_detect_mime(path),
        source_path=str(path),
        file_hash=file_hash(path),
        pages=1,
        full_text=text,
        tokens_estimate=_estimate_tokens(text),
        parsed_at=_now_kst(),
    )


async def _parse_with_kordoc(path: Path) -> ParsedDocument:
    """Call kordoc CLI subprocess for rich format parsing."""
    node = node_bin()
    kordoc = kordoc_cli()
    if not node or not kordoc:
        raise RuntimeError("kordoc runtime not available")

    cmd = [str(node), str(kordoc), str(path), "--format", "json", "--silent"]
    logger.info("Parsing %s via kordoc: %s", path.name, " ".join(cmd))

    def _run() -> subprocess.CompletedProcess:
        return subprocess.run(
            cmd, capture_output=True, text=True, timeout=300,
        )

    result = await asyncio.to_thread(_run)

    if result.returncode != 0:
        stderr = result.stderr[:500] if result.stderr else "unknown error"
        logger.error("kordoc failed for %s: %s", path.name, stderr)
        raise RuntimeError(f"kordoc parse failed: {stderr}")

    try:
        raw = json.loads(result.stdout)
    except json.JSONDecodeError as e:
        raise RuntimeError(f"kordoc returned invalid JSON: {e}")

    markdown = raw.get("markdown") or raw.get("text") or ""
    blocks = raw.get("blocks") or []
    warnings = raw.get("warnings") or []
    metadata = raw.get("metadata") or {}

    sections: list[ParsedSection] = []
    tables: list[ParsedTable] = []
    for block in blocks:
        btype = block.get("type", "")
        if btype == "heading":
            sections.append(ParsedSection(
                title=block.get("text", ""),
                heading_level=block.get("level", 1),
                page=block.get("pageNumber"),
            ))
        elif btype == "table":
            rows = block.get("rows") or []
            cols = len(rows[0]) if rows else 0
            tables.append(ParsedTable(
                page=block.get("pageNumber"),
                rows=len(rows),
                cols=cols,
            ))

    page_count = metadata.get("pageCount") or metadata.get("pages") or 0

    return ParsedDocument(
        doc_id=_new_id(),
        filename=path.name,
        mime=_detect_mime(path),
        source_path=str(path),
        file_hash=file_hash(path),
        pages=page_count,
        full_text=markdown,
        sections=sections,
        tables=tables,
        tokens_estimate=_estimate_tokens(markdown),
        parsed_at=_now_kst(),
        warnings=[str(w) for w in warnings],
    )


async def _ocr_pdf(path: Path) -> ParsedDocument:
    """OCR a scanned PDF — render pages as images, run VLM text extraction."""
    def _extract() -> str:
        import fitz  # type: ignore[import-untyped] — pymupdf
        doc = fitz.open(str(path))
        pages_text: list[str] = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")

            import base64
            b64 = base64.b64encode(img_bytes).decode()
            data_url = f"data:image/png;base64,{b64}"

            from ovo_sidecar.mlx_vlm_runner import vlm_runner
            import asyncio
            loop = asyncio.new_event_loop()
            try:
                result = loop.run_until_complete(
                    vlm_runner.generate(
                        prompt="이 이미지의 모든 텍스트를 정확히 읽어서 그대로 출력해주세요. 다른 설명 없이 텍스트만.",
                        images=[data_url],
                        max_tokens=2048,
                    )
                )
                pages_text.append(f"# Page {i + 1}\n\n{result.text}")
            except Exception as e:
                logger.warning("OCR failed for page %d: %s", i + 1, e)
                pages_text.append(f"# Page {i + 1}\n\n[OCR failed]")
            finally:
                loop.close()
        doc.close()
        return "\n\n".join(pages_text)

    full_text = await asyncio.to_thread(_extract)
    return ParsedDocument(
        doc_id=_new_id(),
        filename=path.name,
        mime="application/pdf",
        source_path=str(path),
        file_hash=file_hash(path),
        pages=full_text.count("# Page "),
        full_text=full_text,
        tokens_estimate=_estimate_tokens(full_text),
        parsed_at=_now_kst(),
        warnings=["OCR-extracted"],
    )


async def parse_bytes(data: bytes, filename: str) -> ParsedDocument:
    """Parse from in-memory bytes (e.g. file upload). Writes temp file then parses."""
    import tempfile

    ext = Path(filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported format: {ext}")

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = Path(tmp.name)

    try:
        doc = await parse_file(tmp_path)
        doc.filename = filename
        return doc
    finally:
        tmp_path.unlink(missing_ok=True)
